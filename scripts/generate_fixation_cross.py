"""Generate a Thaler et al. (2013) ABC optimal fixation target as PowerPoint.

Reference:
    Thaler L, Schütz AC, Goodale MA, Gegenfurtner KR (2013).
    "What is the best fixation target? The effect of target shape on
    stability of fixational eye movements." Vision Research 76:31-42.
    https://doi.org/10.1016/j.visres.2012.10.012

The ABC target = outer filled disk + crosshair cutout + inner filled disk.
Empirically produces the lowest fixation variability of the tested targets.

Geometry is computed for the Pitt BOLDscreen setup:
    - Viewing distance: 139 cm (eye to screen via mirror)
    - Screen: 69.84 cm x 39.29 cm, 1920 x 1080 px
    - Field of view: 28.20 deg x 16.09 deg
    - Resolution: ~68.09 px/deg

Run:
    uv run scripts/generate_fixation_cross.py
"""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

# --- Display geometry (Pitt BOLDscreen, average subject) ----------------------
VIEWING_DISTANCE_CM = 139.0
SCREEN_W_CM, SCREEN_H_CM = 69.84, 39.29
SCREEN_W_PX, SCREEN_H_PX = 1920, 1080


def cm_to_deg(size_cm: float, dist_cm: float) -> float:
    """Visual angle subtended by `size_cm` at `dist_cm` viewing distance."""
    return math.degrees(2 * math.atan(size_cm / (2 * dist_cm)))


FOV_W_DEG = cm_to_deg(SCREEN_W_CM, VIEWING_DISTANCE_CM)  # ~28.20
FOV_H_DEG = cm_to_deg(SCREEN_H_CM, VIEWING_DISTANCE_CM)  # ~16.09
PX_PER_DEG = SCREEN_W_PX / FOV_W_DEG                      # ~68.09

# --- Thaler ABC target parameters ---------------------------------------------
# Thaler et al. 2013 specify outer RADIUS = 0.6°, so outer DIAMETER = 1.2°.
# Inner disk and crosshair preserve Thaler's 3:1 and 6:1 ratios.
OUTER_DEG = 1.2       # outer disk diameter
INNER_DEG = 0.4       # inner disk diameter
CROSS_DEG = 0.2       # crosshair thickness

# --- Colors --------------------------------------------------------------------
BG_GRAY = RGBColor(128, 128, 128)   # iso-luminant mid-gray
FG_BLACK = RGBColor(0, 0, 0)        # maximum contrast for the target

# --- PPT geometry --------------------------------------------------------------
# Slide exactly 1920 x 1080 px at 6350 EMU/px (no projector rescaling).
EMU_PER_PX = 6350
SLIDE_W_EMU = SCREEN_W_PX * EMU_PER_PX   # 12_192_000
SLIDE_H_EMU = SCREEN_H_PX * EMU_PER_PX   # 6_858_000


def deg_to_emu(deg: float) -> int:
    return round(deg * PX_PER_DEG * EMU_PER_PX)


def add_centered_shape(slide, shape_id, width_emu, height_emu, fill_color):
    cx = SLIDE_W_EMU // 2
    cy = SLIDE_H_EMU // 2
    shape = slide.shapes.add_shape(
        shape_id,
        left=Emu(cx - width_emu // 2),
        top=Emu(cy - height_emu // 2),
        width=Emu(width_emu),
        height=Emu(height_emu),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no outline; fill only
    shape.shadow.inherit = False
    return shape


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)

    # --- Slide 1: the fixation target ----------------------------------------
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # Mid-gray background (full-slide rectangle, since changing slide master
    # bg is inconsistent across PPT viewers).
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        Emu(SLIDE_W_EMU), Emu(SLIDE_H_EMU),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_GRAY
    bg.line.fill.background()

    outer_emu = deg_to_emu(OUTER_DEG)
    inner_emu = deg_to_emu(INNER_DEG)
    cross_emu = deg_to_emu(CROSS_DEG)

    # (1) Outer black disk
    add_centered_shape(slide, MSO_SHAPE.OVAL, outer_emu, outer_emu, FG_BLACK)
    # (2) Horizontal gray bar (cuts disk)
    add_centered_shape(slide, MSO_SHAPE.RECTANGLE, outer_emu, cross_emu, BG_GRAY)
    # (3) Vertical gray bar (cuts disk)
    add_centered_shape(slide, MSO_SHAPE.RECTANGLE, cross_emu, outer_emu, BG_GRAY)
    # (4) Inner black disk
    add_centered_shape(slide, MSO_SHAPE.OVAL, inner_emu, inner_emu, FG_BLACK)

    # --- Slide 2: documentation / geometry record ----------------------------
    notes = prs.slides.add_slide(blank_layout)
    notes_bg = notes.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        Emu(SLIDE_W_EMU), Emu(SLIDE_H_EMU),
    )
    notes_bg.fill.solid()
    notes_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
    notes_bg.line.fill.background()

    tb = notes.shapes.add_textbox(
        Emu(600_000), Emu(500_000),
        Emu(SLIDE_W_EMU - 1_200_000), Emu(SLIDE_H_EMU - 1_000_000),
    )
    tf = tb.text_frame
    tf.word_wrap = True

    def add_line(text, size=18, bold=False):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = "Arial"

    add_line("FixationCross.pptx — Thaler ABC optimal fixation target", 24, True)
    add_line("", 10)
    add_line("Design reference:", 16, True)
    add_line(
        "Thaler, Schütz, Goodale & Gegenfurtner (2013). What is the best fixation "
        "target? The effect of target shape on stability of fixational eye "
        "movements. Vision Research 76:31–42.", 14
    )
    add_line("", 10)
    add_line("Display geometry (Pitt BOLDscreen):", 16, True)
    add_line(f"  Viewing distance:   {VIEWING_DISTANCE_CM:.1f} cm", 14)
    add_line(f"  Screen size:        {SCREEN_W_CM} × {SCREEN_H_CM} cm", 14)
    add_line(f"  Native resolution:  {SCREEN_W_PX} × {SCREEN_H_PX} px", 14)
    add_line(f"  Field of view:      {FOV_W_DEG:.2f}° × {FOV_H_DEG:.2f}°", 14)
    add_line(f"  Scale:              {PX_PER_DEG:.2f} px/° "
             f"({10000/PX_PER_DEG:.1f} µm/px on screen)", 14)
    add_line("", 10)
    add_line("Target parameters:", 16, True)
    add_line(f"  Outer disk diameter:    {OUTER_DEG}°  "
             f"({OUTER_DEG * PX_PER_DEG:.1f} px, "
             f"{OUTER_DEG * PX_PER_DEG / PX_PER_DEG * SCREEN_W_CM / FOV_W_DEG:.2f} cm)", 14)
    add_line(f"  Inner disk diameter:    {INNER_DEG}°  "
             f"({INNER_DEG * PX_PER_DEG:.1f} px)", 14)
    add_line(f"  Crosshair thickness:    {CROSS_DEG}°  "
             f"({CROSS_DEG * PX_PER_DEG:.1f} px)", 14)
    add_line(f"  Foreground:             RGB(0, 0, 0)   (black)", 14)
    add_line(f"  Background:             RGB(128, 128, 128)   (mid-gray, "
             f"iso-luminant with typical task stimuli)", 14)
    add_line("", 10)
    add_line("Usage notes:", 16, True)
    add_line(
        "  • Present full-screen on the BOLDscreen. Slide is authored at exactly "
        "1920 × 1080 px (6350 EMU/px) so PowerPoint coordinates map 1:1 to "
        "display pixels — no rescaling.", 13
    )
    add_line(
        "  • Visual angles assume 139 cm viewing distance. If a participant's head "
        "is unusually small/large, recompute px/° before publishing.", 13
    )
    add_line(
        "  • For PsychoPy usage, prefer drawing the target programmatically with "
        "deg units (see scripts/generate_fixation_cross.py) rather than importing "
        "this .pptx as an image.", 13
    )
    return prs


def main() -> None:
    out = Path(__file__).resolve().parents[1] / "materials" / "FixationCross.pptx"
    out.parent.mkdir(exist_ok=True)
    build().save(out)
    print(f"Wrote {out}")
    print(f"  {PX_PER_DEG:.2f} px/°  |  "
          f"outer={OUTER_DEG * PX_PER_DEG:.1f}px  "
          f"inner={INNER_DEG * PX_PER_DEG:.1f}px  "
          f"cross={CROSS_DEG * PX_PER_DEG:.1f}px")


if __name__ == "__main__":
    main()
